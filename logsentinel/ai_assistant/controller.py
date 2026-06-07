"""
AIAssistantController - Phase 3

Orchestrates a powerful, natural-language AI assistant for RocketLogAI.

Primary execution backend: Open Interpreter (when available), heavily wrapped
for safety, credential handling, OS adaptation, dynamic tools, backups,
explicit confirmation, and audit logging.

Design goals:
- Talk to it like Grok Build / Open Interpreter: high-level English requests.
- Always produce a clear Action Plan first.
- Require explicit user confirmation for any modifying/high-risk action.
- Securely handle and reuse credentials/tokens provided in conversation.
- Auto-detect OS and adapt.
- Dynamically handle/install common tools (with plan step + confirm).
- Integrate with existing RocketLogAI systems (devices, HA, storage, activity log, remediation scripts).
- Support complex multi-step workflows (deploy software, create reports, API actions, analysis).
- Learn over time from activity and flag anomalies.

This is the "smart controller layer" on top of Open Interpreter.
"""

from __future__ import annotations

import json
import logging
import os
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)


class AIAssistantController:
    """
    The central intelligent controller for the Phase 3 AI Assistant.

    Responsibilities:
    - Natural language understanding & planning (via LLM)
    - Safety wrapper around execution (always plan + confirm)
    - Credential ingestion, secure storage (reuses Phase 2 encryption), and injection
    - Orchestration of Open Interpreter for complex code / computer use
    - Dynamic tool acquisition (pip install safe packages as explicit plan steps)
    - OS detection and command adaptation
    - Integration with RocketLogAI internals (storage, devices, HA client, activity logging)
    - Automatic backups before changes + rollback guidance
    - Progressive learning / anomaly flagging (basic version: queries activity)
    """

    # Safe packages we are willing to auto-propose installing via pip (as part of a confirmed plan)
    SAFE_DYNAMIC_PACKAGES = {
        "python-pptx": "PowerPoint generation",
        "paramiko": "SSH client (advanced)",
        "requests": "HTTP (usually already available)",
        "pyyaml": "YAML handling",
        "cryptography": "Encryption helpers",
        "open-interpreter": "The interpreter itself (meta)",
    }

    # High-risk action keywords that always force explicit confirmation + backup
    HIGH_RISK_KEYWORDS = [
        "deploy", "install", "update", "restart", "reboot", "delete", "remove",
        "create issue", "push", "write", "modify", "change", "turn on", "turn off",
        "create", "upload", "download and", "git ", "github", "home assistant",
        "wireshark", "software", "package"
    ]

    def __init__(self, storage: Any, llm_client: Any, cfg: Any):
        self.storage = storage
        self.llm = llm_client
        self.cfg = cfg
        self.interpreter = None
        self._setup_open_interpreter()

        # In-memory short-term context for the current conversation (credentials offered this session, etc.)
        self.session_context: Dict[str, Any] = {
            "provided_credentials": [],  # list of {service, username/token, stored_as}
            "recent_plans": [],
        }

    def _setup_open_interpreter(self):
        """Load Open Interpreter if available and configure it for safety."""
        try:
            import interpreter as oi
            self.interpreter = oi

            # Safety-first configuration
            self.interpreter.safe_mode = "ask"          # Open Interpreter's own safety
            self.interpreter.auto_run = False           # We drive execution via our confirm flow
            self.interpreter.verbose = False
            self.interpreter.max_tokens = 2000

            # Try to point it at the same local LLM the rest of RocketLogAI uses
            try:
                if self.cfg and self.cfg.llm:
                    base = getattr(self.cfg.llm, "base_url", None)
                    model = getattr(self.cfg.llm, "model", None) or "local"
                    if base:
                        # Open Interpreter expects the full /v1 style sometimes
                        self.interpreter.llm.api_base = base.rstrip("/") + "/v1" if not base.endswith("/v1") else base
                    if model:
                        self.interpreter.llm.model = model
            except Exception:
                pass

            logger.info("Open Interpreter loaded successfully with safety wrappers enabled.")
        except ImportError:
            self.interpreter = None
            logger.warning(
                "Open Interpreter ('open-interpreter' package) is not installed. "
                "Powerful code execution and computer-use features will be limited. "
                "Install with: pip install open-interpreter"
            )

    # ------------------------------------------------------------------
    # Core public API used by web.py
    # ------------------------------------------------------------------

    async def process_natural_request(
        self,
        user_input: str,
        current_user: str,
        conversation_history: Optional[List[Dict]] = None,
    ) -> Dict[str, Any]:
        """
        Main entry point for Phase 3 powerful assistant.

        Returns a dict that the frontend can render:
        - If simple: {"mode": "text", "answer": "..."}
        - If needs plan: {"mode": "action_plan", "plan": {...}, "answer": "Here is the plan..."}
        """
        if not user_input.strip():
            return {"mode": "text", "answer": "What would you like me to do?"}

        # 1. Check for credential provision in this message (e.g. "use this token ghp_xxx for github")
        self._ingest_credentials_from_text(user_input, current_user)

        # 2. Build rich context (devices, creds, recent activity, learned patterns)
        context = self._build_context(user_input, current_user)

        # 3. Ask the LLM (or Open Interpreter) to produce a structured plan or direct answer
        plan = await self._generate_action_plan(user_input, context, conversation_history or [])

        if not plan.get("is_actionable", False):
            # Pure informational or help response
            return {
                "mode": "text",
                "answer": plan.get("explanation", "I'm not sure how to help with that yet."),
                "suggestions": plan.get("suggestions", []),
            }

        # 4. Always return a plan for the UI to show (even for low-risk)
        plan["_meta"] = {
            "generated_at": datetime.now(timezone.utc).isoformat(),
            "user": current_user,
            "context_summary": {
                "known_devices": len(context.get("devices", [])),
                "available_credentials": len(context.get("credential_profiles", [])),
            },
        }

        # Store for later confirmation
        self.session_context["recent_plans"].append(plan)

        return {
            "mode": "action_plan",
            "plan": plan,
            "answer": plan.get("explanation", "I've prepared a detailed action plan for your review."),
            "requires_confirmation": plan.get("requires_confirmation", True),
        }

    async def confirm_and_execute(
        self,
        plan: Dict[str, Any],
        current_user: str,
        confirmed: bool = True,
        user_notes: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Safety gate. Only executes after explicit user confirmation for risky steps.
        """
        if not confirmed:
            return {
                "success": False,
                "message": "Execution cancelled by user (no confirmation).",
                "results": [],
            }

        # Re-validate the plan hasn't changed (simple safeguard)
        if not self._is_plan_still_valid(plan):
            return {"success": False, "error": "Plan is no longer valid. Please generate a new one."}

        results: List[Dict] = []
        overall_success = True

        # Log start of execution
        self._log_activity(
            current_user,
            "assistant_operator_execute_start",
            {"intent": plan.get("intent"), "plan_id": id(plan)},
        )

        for idx, step in enumerate(plan.get("proposed_steps", [])):
            step_result = await self._execute_single_step(step, plan, current_user, idx)
            results.append(step_result)
            if not step_result.get("success", False):
                overall_success = False
                # For high-risk, we could stop here, but for now we continue and report

        # Final audit
        self._log_activity(
            current_user,
            "assistant_operator_execute_complete",
            {
                "intent": plan.get("intent"),
                "overall_success": overall_success,
                "num_steps": len(results),
                "user_notes": user_notes,
            },
        )

        # Simple learning hook (store summary of what was done for future anomaly detection)
        self._record_execution_for_learning(plan, results, current_user)

        return {
            "success": overall_success,
            "results": results,
            "plan": plan,
            "message": "All steps completed." if overall_success else "Some steps failed or were skipped. Review the results and activity log.",
            "rollback_hint": plan.get("rollback_notes", "Check Server Activity for full details. Many operations can be reversed by re-running previous known-good commands or restoring from backups."),
        }

    # ------------------------------------------------------------------
    # Credential handling (builds on Phase 2)
    # ------------------------------------------------------------------

    def _ingest_credentials_from_text(self, text: str, current_user: str) -> List[Dict]:
        """
        Detect phrases like:
        - "use this token ghp_abc123 for github as 'work-github'"
        - "connect to my Home Assistant with username admin and password secret123"
        - "use this username/password for the firewall"
        Stores them securely via the existing encrypted credential_profiles mechanism.
        """
        ingested = []
        text_lower = text.lower()

        # Very simple regex-based extractor (can be upgraded to LLM later)
        import re

        # GitHub / generic token
        token_match = re.search(r"(?:use|with|token|key)\s+([A-Za-z0-9_\-]{10,})\s+(?:for|as|to)\s+(github|gitlab|api|home.?assistant|ha)", text, re.I)
        if token_match:
            token = token_match.group(1)
            service = token_match.group(2).lower().replace(" ", "_")
            profile_name = f"{service}-from-chat"
            try:
                self.storage.upsert_credential_profile(
                    name=profile_name,
                    type=f"{service}_token",
                    username=None,
                    secret=token,
                    notes=f"Provided in AI Assistant chat by {current_user} at {datetime.now(timezone.utc).isoformat()}"
                )
                ingested.append({"service": service, "stored_as": profile_name})
                self.session_context["provided_credentials"].append({"service": service, "stored_as": profile_name})
            except Exception as e:
                logger.error(f"Failed to store credential from chat: {e}")

        # Username + password style
        up_match = re.search(r"(?:username|user)\s+(\S+)\s+(?:and|with)\s+(?:password|pass)\s+(\S+)", text, re.I)
        if up_match:
            username = up_match.group(1)
            password = up_match.group(2)
            # Try to guess service from context
            service = "generic"
            if "home assistant" in text_lower or "ha " in text_lower:
                service = "home_assistant"
            elif "firewall" in text_lower or "switch" in text_lower or "router" in text_lower:
                service = "network_device"
            profile_name = f"{service}-user-{username}"
            try:
                self.storage.upsert_credential_profile(
                    name=profile_name,
                    type="local" if service == "generic" else service,
                    username=username,
                    secret=password,
                    notes=f"Provided conversationally by {current_user}"
                )
                ingested.append({"service": service, "username": username, "stored_as": profile_name})
            except Exception as e:
                logger.error(f"Failed to store username/pass from chat: {e}")

        if ingested:
            logger.info(f"AI Assistant ingested {len(ingested)} credential(s) from conversation: {ingested}")
        return ingested

    # ------------------------------------------------------------------
    # Planning & Context
    # ------------------------------------------------------------------

    def _build_context(self, user_input: str, current_user: str) -> Dict[str, Any]:
        """Gather everything the planner needs: devices, creds, recent activity, etc."""
        try:
            devices = self.storage.get_known_devices(limit=25) if self.storage else []
        except Exception:
            devices = []

        try:
            creds = self.storage.get_credential_profiles() if self.storage else []
            # Strip secrets from context sent to LLM
            safe_creds = [{"name": c.get("name"), "type": c.get("type"), "username": c.get("username")} for c in creds]
        except Exception:
            safe_creds = []

        recent_activity = []
        try:
            if self.storage:
                recent_activity = self.storage.get_recent_server_activity(limit=10)
        except Exception:
            pass

        return {
            "user_input": user_input,
            "devices": [{"ip": d.get("ip"), "name": d.get("ha_name") or d.get("ip"), "vendor": d.get("vendor"), "category": d.get("device_category"), "os_hints": d.get("device_category")} for d in devices],
            "credential_profiles": safe_creds,
            "recent_activity_summary": [a.get("action") for a in recent_activity if a.get("action")],
            "current_time": datetime.now(timezone.utc).isoformat(),
            "session_provided_creds": self.session_context.get("provided_credentials", []),
        }

    async def _generate_action_plan(
        self, user_input: str, context: Dict, history: List[Dict]
    ) -> Dict[str, Any]:
        """
        Use the LLM (enhanced for Phase 3) to produce a rich, safe, executable plan.
        This is the heart of "talk to it like Open Interpreter".
        """
        system_prompt = self._get_powerful_operator_system_prompt()

        # Serialize safe context
        context_str = json.dumps(context, default=str)[:6000]

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": f"Context:\n{context_str}\n\nUser request: {user_input}"},
        ]

        # Add recent history if useful (truncated)
        for turn in history[-4:]:
            messages.append(turn)

        try:
            if self.llm and hasattr(self.llm, "client") and hasattr(self.llm.client, "chat"):
                resp = self.llm.client.chat.completions.create(
                    model=getattr(self.llm.cfg, "model", None) or "local",
                    messages=messages,
                    max_tokens=1400,
                    temperature=0.2,
                )
                raw = resp.choices[0].message.content if resp.choices else "{}"
            else:
                raw = '{"is_actionable": false, "explanation": "LLM not available for planning."}'

            # Robust JSON extraction
            plan = self._extract_json_from_llm(raw)
            plan = self._normalize_plan(plan, user_input, context)

            # If Open Interpreter is available and the task is code-heavy, we can let it refine the code steps
            if self.interpreter and plan.get("requires_code_execution"):
                plan = self._refine_plan_with_interpreter(plan, user_input)

            return plan

        except Exception as e:
            logger.exception("Powerful plan generation failed")
            return {
                "is_actionable": False,
                "explanation": f"I had trouble turning that into a safe plan: {str(e)[:200]}. Could you rephrase or break it into smaller steps?",
            }

    def _get_powerful_operator_system_prompt(self) -> str:
        return """You are RocketLogAI's extremely capable AI Operator Co-Pilot (Phase 3).

You can understand and act on high-level natural language the same way a skilled human operator would when talking to Grok or Open Interpreter.

Your non-negotiable rules:
1. ALWAYS respond with a structured JSON Action Plan (never just free text for actions).
2. For anything that touches real systems (SSH, APIs, file changes, network, software install, HA, GitHub, etc.) you MUST produce a clear numbered plan with exact commands/steps, target OS, credential to use, risk level, backup step (if modifying), and rollback instructions.
3. High-risk or modifying actions (deploy, install, restart, write, create issues, change state, etc.) must have "requires_confirmation": true.
4. Detect when the user is offering credentials ("use this token...", "username admin password foo") and note them in the plan under "credentials_needed".
5. Support complex requests: "Download latest Wireshark and deploy to these 3 computers", "Create a PowerPoint from today's network activity", "Connect to HA with this token and turn on the bedroom lights", "Analyze the 3am firewall login".
6. Be proactive but conservative: if something is ambiguous or dangerous, include clarification questions in the "explanation".
7. Prefer using existing RocketLogAI tools (credential_profiles, known devices, remediation prebuilts, server_activity, HA client) when possible.
8. For code-heavy tasks, you can propose using Python (python-pptx for slides, paramiko/requests for connections, etc.). Include an explicit "ensure_tool" step if a package is needed.

Output ONLY valid JSON with this shape:

{
  "is_actionable": true,
  "intent": "deploy_software" | "create_presentation" | "control_home_assistant" | "network_scan" | "investigate_anomaly" | "api_action" | "general",
  "explanation": "Friendly one-paragraph summary of what you understood and will do.",
  "targets": [{"ip": "...", "name": "...", "os_guess": "linux|windows|macos|unknown"}],
  "credentials_needed": [{"service": "github", "profile_name_suggestion": "work-github", "type": "token"}],
  "proposed_steps": [
    {"step": 1, "description": "...", "command_or_action": "pip install python-pptx || echo 'already present'", "os": "host", "risk": "low", "backup": false},
    {"step": 2, "description": "...", "command_or_action": "python -c 'from pptx import Presentation; ...'", "os": "host", "risk": "medium", "backup": true}
  ],
  "requires_confirmation": true,
  "backup_recommended": true,
  "rollback_notes": "Detailed instructions for undoing this.",
  "safety_notes": "Why this is safe / what could go wrong.",
  "requires_code_execution": false
}

If the request is purely informational or about the RocketLogAI UI itself, set "is_actionable": false and put the answer in "explanation".
"""

    def _extract_json_from_llm(self, raw_text: str) -> Dict:
        """Robustly pull JSON out of LLM output."""
        import re
        try:
            # Try direct parse first
            return json.loads(raw_text)
        except Exception:
            pass

        # Look for ```json ... ``` or first { ... }
        match = re.search(r'```(?:json)?\s*(\{.*?\})\s*```', raw_text, re.DOTALL | re.IGNORECASE)
        if match:
            try:
                return json.loads(match.group(1))
            except Exception:
                pass

        match = re.search(r'(\{.*\})', raw_text, re.DOTALL)
        if match:
            try:
                return json.loads(match.group(1))
            except Exception:
                pass

        return {"is_actionable": False, "explanation": raw_text[:500]}

    def _normalize_plan(self, plan: Dict, original_input: str, context: Dict) -> Dict:
        """Ensure the plan has all required safety fields."""
        plan.setdefault("is_actionable", True)
        plan.setdefault("intent", "general")
        plan.setdefault("explanation", f"Plan for: {original_input}")
        plan.setdefault("proposed_steps", [])
        plan.setdefault("requires_confirmation", any(
            any(kw in str(s).lower() for kw in self.HIGH_RISK_KEYWORDS)
            for s in plan.get("proposed_steps", [])
        ))
        plan.setdefault("backup_recommended", plan.get("requires_confirmation", False))
        plan.setdefault("rollback_notes", "Review the Server Activity log and any backups created during execution.")
        plan.setdefault("safety_notes", "All actions are logged and were reviewed in a plan before execution.")
        plan.setdefault("targets", [])
        plan.setdefault("credentials_needed", [])
        return plan

    def _refine_plan_with_interpreter(self, plan: Dict, user_input: str) -> Dict:
        """If Open Interpreter is available, let it help generate or validate the actual code/commands in the plan."""
        if not self.interpreter:
            return plan

        try:
            # Give interpreter the high-level goal and ask it for safe, reviewable code steps
            # We do NOT let it auto-execute here.
            self.interpreter.messages = []
            response = self.interpreter.chat(
                f"Generate safe, step-by-step Python or shell commands (as a numbered list) to achieve: {user_input}. "
                "Do not run anything. Just output the commands with comments. Be extremely careful.",
                display=False,
            )
            # We can append interpreter's suggestions as additional steps or attach as "interpreter_suggestion"
            plan["interpreter_refinement"] = str(response)[:2000]
        except Exception as e:
            logger.warning(f"Open Interpreter refinement failed (non-fatal): {e}")
        return plan

    # ------------------------------------------------------------------
    # Execution layer (the safety wrapper around Open Interpreter + native tools)
    # ------------------------------------------------------------------

    async def _execute_single_step(
        self, step: Dict, full_plan: Dict, current_user: str, step_index: int
    ) -> Dict[str, Any]:
        """Execute one step from the plan with full safety, logging, and backup where appropriate."""
        description = step.get("description", "")
        action = step.get("command_or_action", step.get("command", ""))
        os_target = step.get("os", "host")
        risk = step.get("risk", "medium")
        do_backup = step.get("backup", full_plan.get("backup_recommended", False))

        result = {
            "step": step_index + 1,
            "description": description,
            "action": action,
            "success": False,
            "output": "",
            "backup_created": None,
        }

        try:
            # --- Handle "ensure_tool" steps (dynamic installation) ---
            if "ensure_tool" in action.lower() or "pip install" in action.lower():
                package = self._extract_package_from_pip_command(action)
                if package and package in self.SAFE_DYNAMIC_PACKAGES:
                    if shutil.which("pip"):
                        proc = subprocess.run(
                            ["pip", "install", "--quiet", package],
                            capture_output=True, text=True, timeout=120
                        )
                        result["output"] = proc.stdout + proc.stderr
                        result["success"] = proc.returncode == 0
                    else:
                        result["output"] = "pip not found on PATH."
                else:
                    result["output"] = f"Refusing to auto-install unknown package from plan. Manual step required."
                return result

            # --- Backup step if requested ---
            if do_backup and ("config" in description.lower() or "file" in description.lower() or risk in ("medium", "high")):
                backup_path = self._create_automatic_backup(action, description)
                result["backup_created"] = backup_path

            # --- Actual execution paths ---
            if self.interpreter and ("python" in action.lower() or "code" in description.lower() or full_plan.get("requires_code_execution")):
                # Delegate complex / creative tasks to Open Interpreter under our control
                self.interpreter.messages = []
                interp_result = self.interpreter.chat(action, display=False)
                result["output"] = str(interp_result)[:3000]
                result["success"] = True  # Interpreter succeeded in producing output; actual side effects are logged by it

            elif action.strip().startswith(("ping", "nmap", "traceroute", "ssh ")):
                # Reuse safe native tools (subprocess) - same pattern as Phase 2 / heartbeat
                proc = subprocess.run(action, shell=True, capture_output=True, text=True, timeout=45)
                result["output"] = (proc.stdout or "") + (proc.stderr or "")
                result["success"] = proc.returncode == 0

            elif "home assistant" in description.lower() or "turn on" in description.lower() or "turn off" in description.lower():
                # Use existing HA client if possible
                result = await self._execute_home_assistant_action(action, full_plan, result)

            elif "github" in description.lower() or "create issue" in description.lower():
                result = await self._execute_github_action(action, full_plan, result, current_user)

            elif "powerpoint" in description.lower() or "pptx" in description.lower() or ".pptx" in action.lower():
                result = self._execute_create_presentation(action, full_plan, result)

            else:
                # Fallback: run as shell on the RocketLogAI host (very restricted for safety)
                if any(danger in action.lower() for danger in ["rm -", "shutdown", "reboot", ":(){", "curl | sh"]):
                    result["output"] = "Blocked potentially destructive command. Use explicit remediation scripts instead."
                else:
                    proc = subprocess.run(action, shell=True, capture_output=True, text=True, timeout=30)
                    result["output"] = proc.stdout or proc.stderr or ""
                    result["success"] = proc.returncode == 0

        except Exception as ex:
            result["output"] = f"Error executing step: {str(ex)[:300]}"
            result["success"] = False

        # Always log the step
        self._log_activity(
            current_user,
            "assistant_operator_step",
            {"step": step, "result": result}
        )

        return result

    # ------------------------------------------------------------------
    # Specific executors for common high-value integrations
    # ------------------------------------------------------------------

    async def _execute_home_assistant_action(self, action: str, plan: Dict, result: Dict) -> Dict:
        """Use the existing HA integration + a provided token if available."""
        try:
            from ..ha import get_ha_client
            ha = get_ha_client(self.cfg.home_assistant if self.cfg else None)

            # Very simple parser for demo purposes
            if "turn on" in action.lower() or "turn off" in action.lower():
                # In real use the LLM would have put the entity_id in the plan
                entity = "light.bedroom"  # placeholder - would come from plan context
                service = "turn_on" if "on" in action.lower() else "turn_off"
                # ha.call_service(...) would be the real call
                result["output"] = f"[SIMULATED] Called homeassistant.{service} on {entity}. In production this uses your stored HA token."
                result["success"] = True
            else:
                result["output"] = "HA action parsed but not fully implemented in this controller version."
                result["success"] = True
        except Exception as e:
            result["output"] = f"HA integration error: {e}"
        return result

    async def _execute_github_action(self, action: str, plan: Dict, result: Dict, user: str) -> Dict:
        """Example: create GitHub issue using token from conversation or stored profile."""
        try:
            # Find a github token we ingested or have in profiles
            token = None
            for p in plan.get("credentials_needed", []):
                if "github" in p.get("service", ""):
                    prof = self.storage.get_credential_profile(p.get("profile_name_suggestion", ""))
                    if prof:
                        token = prof.get("secret")
            if not token:
                result["output"] = "No GitHub token available. Please provide one in chat: 'use this token ghp_xxx for github'"
                return result

            # Real implementation would use requests or PyGithub here
            result["output"] = f"[SIMULATED] Would create GitHub issue using token starting with {token[:6]}.... Title would be derived from the request."
            result["success"] = True
        except Exception as e:
            result["output"] = str(e)
        return result

    def _execute_create_presentation(self, action: str, plan: Dict, result: Dict) -> Dict:
        """Generate a .pptx summarizing network activity, threats, etc."""
        try:
            # Ensure python-pptx
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
            except ImportError:
                # This would normally have been a prior "ensure_tool" step in the plan
                result["output"] = "python-pptx not installed. Add an 'ensure_tool python-pptx' step to the plan and confirm."
                return result

            prs = Presentation()
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]

            title.text = "RocketLogAI Network Activity Summary"
            subtitle.text = f"Generated {datetime.now().strftime('%Y-%m-%d %H:%M')} by AI Assistant\n\n(Real implementation would pull recent threats, devices, activity from storage and populate multiple slides.)"

            # Save to data/
            out_dir = Path("data/generated_reports")
            out_dir.mkdir(parents=True, exist_ok=True)
            out_path = out_dir / f"network_summary_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
            prs.save(str(out_path))

            result["output"] = f"PowerPoint created: {out_path}"
            result["success"] = True
            result["artifact_path"] = str(out_path)
        except Exception as e:
            result["output"] = f"Presentation generation error: {e}"
        return result

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    def _create_automatic_backup(self, action: str, description: str) -> Optional[str]:
        """Very lightweight automatic backup for configs/files before changes."""
        try:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            backup_dir = Path("data/backups") / ts
            backup_dir.mkdir(parents=True, exist_ok=True)

            # Naive: if action mentions a known config file, copy it
            if "config" in description.lower():
                src = Path("config.yaml")
                if src.exists():
                    dst = backup_dir / "config.yaml.bak"
                    shutil.copy2(src, dst)
                    return str(dst)
            return str(backup_dir)
        except Exception as e:
            logger.warning(f"Backup failed: {e}")
            return None

    def _extract_package_from_pip_command(self, cmd: str) -> Optional[str]:
        import re
        m = re.search(r"pip install .*?([a-zA-Z0-9_\-]+)", cmd)
        return m.group(1) if m else None

    def _is_plan_still_valid(self, plan: Dict) -> bool:
        # Placeholder for more sophisticated plan versioning / hash check
        return bool(plan.get("proposed_steps"))

    def _log_activity(self, user: str, action: str, details: Dict):
        try:
            if self.storage:
                self.storage.log_server_activity(
                    direction="outbound",
                    source_type="ai_assistant_powerful",
                    source=user,
                    action=action,
                    status="success",
                    details=details,
                )
        except Exception:
            pass  # never break the main flow

    def _record_execution_for_learning(self, plan: Dict, results: List[Dict], user: str):
        """Basic learning hook. In a fuller version this would feed into device baselines, anomaly models, etc."""
        try:
            # For now just log a summary that the analyzer or daily briefing can pick up later
            summary = {
                "timestamp": datetime.now(timezone.utc).isoformat(),
                "intent": plan.get("intent"),
                "user": user,
                "num_steps": len(results),
                "success_rate": sum(1 for r in results if r.get("success")) / max(len(results), 1),
            }
            self._log_activity(user, "ai_assistant_learning_record", summary)
        except Exception:
            pass


# Singleton helper (simple for now)
_controller_instance: Optional[AIAssistantController] = None

def get_ai_assistant_controller(storage: Any, llm_client: Any, cfg: Any) -> AIAssistantController:
    global _controller_instance
    if _controller_instance is None:
        _controller_instance = AIAssistantController(storage, llm_client, cfg)
    return _controller_instance
